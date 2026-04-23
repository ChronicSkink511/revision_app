from __future__ import annotations

import logging
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from revision_app.schemas import DocumentContent

TEXT_EXTS = {".txt", ".md"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}
MAX_TEXT_CHARS = 30000



def _truncate(text: str) -> str:
    return text[:MAX_TEXT_CHARS] if len(text) > MAX_TEXT_CHARS else text



def _extract_pdf(path: Path, max_chars: int = MAX_TEXT_CHARS) -> str:
    reader = PdfReader(str(path))
    chunks: list[str] = []
    total = 0
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if not page_text:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        chunk = page_text[:remaining]
        chunks.append(chunk)
        total += len(chunk)
    return "\n".join(chunks)



def _extract_docx(path: Path, max_chars: int = MAX_TEXT_CHARS) -> str:
    doc = Document(str(path))
    chunks: list[str] = []
    total = 0
    for paragraph in doc.paragraphs:
        text = paragraph.text
        if not text:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        chunk = text[:remaining]
        chunks.append(chunk)
        total += len(chunk)
    return "\n".join(chunks)



def _extract_pptx(path: Path, max_chars: int = MAX_TEXT_CHARS) -> str:
    prs = Presentation(str(path))
    text_items: list[str] = []
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                remaining = max_chars - total
                if remaining <= 0:
                    return "\n".join(text_items)
                chunk = shape.text[:remaining]
                text_items.append(chunk)
                total += len(chunk)
    return "\n".join(text_items)



def _extract_plain(path: Path, max_chars: int = MAX_TEXT_CHARS) -> str:
    with path.open("r", encoding="utf-8", errors="ignore") as handle:
        return handle.read(max_chars)



def parse_documents(file_paths: list[Path], image_analyzer, logger: logging.Logger) -> tuple[list[DocumentContent], list[str]]:
    documents: list[DocumentContent] = []
    warnings: list[str] = []

    for file_path in file_paths:
        suffix = file_path.suffix.lower()
        text = ""
        image_summaries: list[str] = []

        try:
            if suffix == ".pdf":
                text = _extract_pdf(file_path, max_chars=MAX_TEXT_CHARS)
            elif suffix == ".docx":
                text = _extract_docx(file_path, max_chars=MAX_TEXT_CHARS)
            elif suffix == ".pptx":
                text = _extract_pptx(file_path, max_chars=MAX_TEXT_CHARS)
            elif suffix in TEXT_EXTS:
                text = _extract_plain(file_path, max_chars=MAX_TEXT_CHARS)
            elif suffix in IMAGE_EXTS:
                vision = image_analyzer.analyze(file_path)
                image_summaries.append(vision)
                text = f"[Image Analysis]\n{vision}"
            else:
                warnings.append(f"Parser skipped unsupported file: {file_path.name}")
                continue

            content = DocumentContent(
                source_path=file_path,
                file_type=suffix.lstrip("."),
                text=_truncate(text.strip()),
                image_summaries=image_summaries,
            )

            if not content.text:
                warnings.append(f"No extractable content in {file_path.name}")
                continue

            documents.append(content)
        except Exception as exc:
            logger.exception("Parsing failed for %s", file_path)
            warnings.append(f"Corrupted or unreadable file skipped: {file_path.name} ({exc})")

    return documents, warnings
